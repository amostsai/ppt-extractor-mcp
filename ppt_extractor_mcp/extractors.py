"""
PowerPoint 文字提取的核心實現
"""
from typing import List, Dict, Any, BinaryIO, Union
import tempfile
import os
from pptx import Presentation

def extract_text_from_ppt(
    content: Union[bytes, BinaryIO], 
    include_notes: bool = False,
    include_empty_slides: bool = False
) -> List[Dict[str, Any]]:
    """
    從 PowerPoint 檔案中提取文字內容
    
    Args:
        content: PowerPoint 檔案的二進制內容或檔案對象
        include_notes: 是否包含講者備註
        include_empty_slides: 是否包含沒有文字的幻燈片
        
    Returns:
        包含每張幻燈片內容的列表
    """
    # 使用臨時檔案處理 PowerPoint
    with tempfile.NamedTemporaryFile(suffix='.pptx', delete=False) as temp_file:
        temp_file_path = temp_file.name
        
        # 如果是二進制內容，直接寫入
        if isinstance(content, bytes):
            temp_file.write(content)
        # 如果是檔案對象，讀取後寫入
        else:
            temp_file.write(content.read())
    
    try:
        # 打開 PowerPoint 檔案
        prs = Presentation(temp_file_path)
        slides_content = []
        
        for i, slide in enumerate(prs.slides, 1):
            slide_data = {
                "slide_number": i,
                "content": [],
                "notes": ""
            }
            
            # 提取每個形狀中的文字
            for shape in slide.shapes:
                if hasattr(shape, "text") and shape.text.strip():
                    slide_data["content"].append(shape.text.strip())
            
            # 提取講者備註
            if include_notes and slide.has_notes_slide:
                notes_slide = slide.notes_slide
                for shape in notes_slide.shapes:
                    if hasattr(shape, "text") and shape.text.strip():
                        slide_data["notes"] += shape.text.strip() + "\n"
            
            # 只有當幻燈片有內容或用戶選擇包含空幻燈片時才添加
            if include_empty_slides or slide_data["content"] or slide_data["notes"]:
                slides_content.append(slide_data)
        
        return slides_content
        
    finally:
        # 刪除臨時檔案
        if os.path.exists(temp_file_path):
            os.unlink(temp_file_path)

def format_as_markdown(slides_content: List[Dict[str, Any]], filename: str) -> str:
    """
    將提取的內容格式化為 Markdown
    
    Args:
        slides_content: 從 extract_text_from_ppt 返回的幻燈片內容
        filename: 原始檔案名稱
        
    Returns:
        格式化的 Markdown 文本
    """
    if not slides_content:
        return f"# {filename}\n\n*PowerPoint 中沒有找到文字內容。*"
    
    result = [f"# {filename} 內容摘要\n"]
    
    for slide in slides_content:
        slide_num = slide["slide_number"]
        result.append(f"## 幻燈片 {slide_num}\n")
        
        # 處理主要內容
        for text in slide["content"]:
            # 檢查是否為標題文字 (通常較短且在幻燈片頂部)
            if len(text) < 100 and slide["content"].index(text) == 0:
                result.append(f"### {text}\n")
            else:
                result.append(f"{text}\n")
        
        # 處理講者備註
        if slide["notes"]:
            result.append("\n> **講者備註：**\n")
            for line in slide["notes"].split("\n"):
                if line.strip():
                    result.append(f"> {line}\n")
        
        result.append("")  # 添加空行分隔每張幻燈片
    
    return "\n".join(result)

def format_as_plain_text(slides_content: List[Dict[str, Any]], filename: str) -> str:
    """
    將提取的內容格式化為純文字
    
    Args:
        slides_content: 從 extract_text_from_ppt 返回的幻燈片內容
        filename: 原始檔案名稱
        
    Returns:
        格式化的純文字
    """
    if not slides_content:
        return f"{filename}\n\nPowerPoint 中沒有找到文字內容。"
    
    result = [f"{filename} 內容摘要\n"]
    
    for slide in slides_content:
        slide_num = slide["slide_number"]
        result.append(f"[幻燈片 {slide_num}]")
        
        # 處理主要內容
        for text in slide["content"]:
            result.append(text)
        
        # 處理講者備註
        if slide["notes"]:
            result.append("\n講者備註:")
            for line in slide["notes"].split("\n"):
                if line.strip():
                    result.append(f"  {line}")
        
        result.append("-" * 40)  # 添加分隔線
    
    return "\n".join(result)